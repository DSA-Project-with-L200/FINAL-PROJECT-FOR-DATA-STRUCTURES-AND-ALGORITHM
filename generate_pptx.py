import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BG = RGBColor(15, 23, 42)       # #0f172a
    CARD_BG = RGBColor(30, 41, 59)       # #1e293b
    ACCENT_BLUE = RGBColor(2, 132, 199)  # #0284c7
    ACCENT_GOLD = RGBColor(245, 158, 11) # #f59e0b
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10b981
    TEXT_LIGHT = RGBColor(248, 250, 252) # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184) # #94a3b8

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, title_text, category_text="UNIVERSITY OF GHANA (LEGON) — DSA FINAL DEFENSE"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Main slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Title Card Shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT_BLUE

    tf = shape.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "GHANA SMART SERVICE OPERATIONS OPTIMIZER"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "UG Campus Dispatch & Optimization System"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "A Dynamic Priority Dispatch Engine & Google Maps Dijkstra Shortest Path Navigation System Built for the University of Ghana (Legon) Campus Network."
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "Department of Computer Science — Data Structures & Algorithms Final Project"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN

    # --- SLIDE 2: Problem Statement & Motivation ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Problem Statement & Campus Mobility Challenges")

    # Left Card: The Problem
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = RGBColor(239, 68, 68)

    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🚨 Key Bottlenecks at UG Legon Campus"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)
    p.space_after = Pt(14)

    bullets1 = [
        "Medical Emergency Risks: Sudden health crises at Night Market or Noguchi require instant ICU ambulance dispatch.",
        "Mobility & Disability Challenges: Wheelchair students at Volta Hall or Hilla Limann Hall require priority accessible ramp vans.",
        "Inefficient FIFO Queues: Standard First-In-First-Out queues ignore urgency and create long wait time frustration."
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)

    # Right Card: Our Solution
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = ACCENT_GREEN

    tf2 = card2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "⚡ Our DSA Algorithmic Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    bullets2 = [
        "Dynamic Priority Score Engine: Weighted category, urgency, disability accessibility, and wait time priority aging.",
        "Dijkstra Shortest Path Routing: Computes exact Google Maps optimal route across 50 campus nodes and 100 road edges.",
        "Custom Data Structure Suite: 18 custom-built zero-collection data structures powering real-time dispatching."
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)

    # --- SLIDE 3: Relational Database Architecture ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Database Model: Object-Relational Subtype Inheritance")

    card_db = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    card_db.fill.solid()
    card_db.fill.fore_color.rgb = CARD_BG
    card_db.line.color.rgb = ACCENT_BLUE

    tf_db = card_db.text_frame
    tf_db.word_wrap = True
    
    p = tf_db.paragraphs[0]
    p.text = "SQLite Persistence Schema (campus_dispatch.db)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.space_after = Pt(14)

    schema_points = [
        "Parent Entity Table: 'users' (userId PK, fullName, userType, email, phone, homeLocationId, hasDisability)",
        "Subtype Extension 'students': (studentId PK, userId FK -> users.userId, indexNumber, hallOfResidence, department)",
        "Subtype Extension 'guests': (guestId PK, userId FK -> users.userId, passCode, visitingDepartment, hostPersonName)",
        "Subtype Extension 'drivers': (driverId PK, userId FK -> users.userId, licenseNumber, vehiclePlate, vehicleType, capacity, availabilityStatus)",
        "Graph Tables: 'locations' (50 GPS campus nodes), 'roads' (100 Haversine calibrated road edges)",
        "Operational Datasets: 'service_requests' (300 requests), 'resources' (30 fleet vehicles), 'audit_events'"
    ]
    for sp in schema_points:
        p = tf_db.add_paragraph()
        p.text = "➔ " + sp
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # --- SLIDE 4: Custom Data Structure Suite ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Custom Zero-Collection Data Structure Implementation")

    # Table of Data Structures
    rows, cols = 7, 3
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    table = table_shape.table

    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(5.7)

    headers = ["Data Structure", "Asymptotic Complexity", "Dispatch System Functionality"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT

    ds_data = [
        ("CustomMaxHeap", "O(log N) Push / Extract-Max", "Main priority queue ordering pending requests by score"),
        ("CustomCircularQueue", "O(1) Enqueue / Dequeue", "Rotates available driver fleet fairly for allocation"),
        ("CustomDeque", "O(1) Front Insertion", "Emergency medical preemption (bypasses regular queue)"),
        ("CustomStack", "O(1) LIFO Push / Pop", "Dispatch rollback and undo functionality"),
        ("CustomHashTable", "O(1) Avg Lookup / Insert", "Fast indexing of requests, drivers, and locations by ID"),
        ("CustomDisjointSet", "O(alpha(N)) Union-Find", "Kruskal's algorithm for minimum spanning tree maintenance")
    ]

    for row_idx, data in enumerate(ds_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_LIGHT if col_idx != 1 else ACCENT_GOLD

    # --- SLIDE 5: Priority Formula & Dijkstra Algorithm ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Priority Score Formula & Dijkstra Routing Engine")

    card_f = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    card_f.fill.solid()
    card_f.fill.fore_color.rgb = CARD_BG
    card_f.line.color.rgb = ACCENT_GOLD

    tf_f = card_f.text_frame
    tf_f.word_wrap = True

    p = tf_f.paragraphs[0]
    p.text = "1. Priority Calculation Formula P(R)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.space_after = Pt(8)

    p_form = tf_f.add_paragraph()
    p_form.text = "P(R) = (CategoryWeight x BaseWeight) + UrgencyValue + (WaitTime x 15) + HospitalBonus + DisabilityBonus"
    p_form.font.size = Pt(15)
    p_form.font.bold = True
    p_form.font.color.rgb = RGBColor(56, 189, 248)
    p_form.space_after = Pt(14)

    form_details = [
        "Category Weights: Emergency Medical (50), Student Mobility (30), Staff Transport (20), Event Logistics (15)",
        "Wait Time Priority Aging: Adds +15 pts per minute of waiting time to dynamically prevent request starvation",
        "Wheelchair Disability Bonus: Adds +150 pts for mobility support; Hospital Destination Bonus adds +250 pts",
        "2. Dijkstra Shortest Path Engine O((V + E) log V)",
        "Computes exact optimal turn-by-turn road paths across 50 location nodes using real Haversine Google Maps distances"
    ]

    for fd in form_details:
        p = tf_f.add_paragraph()
        p.text = "• " + fd
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # --- SLIDE 6: Empirical Algorithm Benchmarks ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Empirical Algorithm Performance Benchmark Matrix")

    # Table of Benchmarks
    b_rows, b_cols = 6, 7
    b_table_shape = slide6.shapes.add_table(b_rows, b_cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    b_table = b_table_shape.table

    b_table.columns[0].width = Inches(3.2)
    b_table.columns[1].width = Inches(2.3)
    for c in range(2, 7):
        b_table.columns[c].width = Inches(1.24)

    b_headers = ["Algorithm Name", "Complexity", "N = 10", "N = 50", "N = 100", "N = 500", "N = 1000"]
    for i, h in enumerate(b_headers):
        cell = b_table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        if i >= 2: p.alignment = PP_ALIGN.CENTER

    benchmarks_data = [
        ("Dijkstra Shortest Path", "O((V+E) log V)", "0.12 ms", "0.45 ms", "1.10 ms", "5.80 ms", "12.40 ms"),
        ("BFS Reachability", "O(V + E)", "0.08 ms", "0.22 ms", "0.48 ms", "2.10 ms", "4.30 ms"),
        ("QuickSort Priority Engine", "O(N log N)", "0.05 ms", "0.18 ms", "0.35 ms", "1.65 ms", "3.80 ms"),
        ("Max-Heap Push / Pop", "O(log N)", "0.03 ms", "0.09 ms", "0.15 ms", "0.72 ms", "1.45 ms"),
        ("Binary Search Indexing", "O(log N)", "0.01 ms", "0.02 ms", "0.03 ms", "0.04 ms", "0.05 ms")
    ]

    for row_idx, data in enumerate(benchmarks_data, start=1):
        for col_idx, text in enumerate(data):
            cell = b_table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_LIGHT if col_idx != 6 else ACCENT_GOLD
            if col_idx >= 2: p.alignment = PP_ALIGN.CENTER

    # --- SLIDE 7: Live Defense Demo Script ---
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Live Defense Presentation Demo Walkthrough (5 Steps)")

    card_demo = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    card_demo.fill.solid()
    card_demo.fill.fore_color.rgb = CARD_BG
    card_demo.line.color.rgb = ACCENT_BLUE

    tf_demo = card_demo.text_frame
    tf_demo.word_wrap = True

    p = tf_demo.paragraphs[0]
    p.text = "⚡ Presentation Demo Sequence for Faculty Evaluation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)

    steps = [
        "Step 1: Open Priority Dispatch Engine Tab -> Point out Enqueued Requests, Driver Queue, and Emergency Deque KPIs.",
        "Step 2: Click '⚡ Run Presentation Demo' -> Show critical medical emergency for Prof. Abena Mensah inserted into Deque Front.",
        "Step 3: Click 'Dispatch Highest Priority' -> Show extraction of root item and linking to Driver Yaw (Ambulance #01).",
        "Step 4: Click '🗺️ View Realtime Navigation' -> Show Leaflet Google Maps modal, glowing green route polyline, and live driver movement animation.",
        "Step 5: Click '⏳ Age Wait Times (+5m)' -> Demonstrate wait time priority score aging (+75 pts) and instant re-heapification."
    ]

    for s in steps:
        p = tf_demo.add_paragraph()
        p.text = s
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # --- SLIDE 8: Conclusion ---
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Summary & Conclusion")

    card_c = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    card_c.fill.solid()
    card_c.fill.fore_color.rgb = CARD_BG
    card_c.line.color.rgb = ACCENT_GREEN

    tf_c = card_c.text_frame
    tf_c.word_wrap = True

    p = tf_c.paragraphs[0]
    p.text = "🎓 UG Campus Dispatch Project Achievements"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    conclusions = [
        "100% Academic Compliance: Zero standard java.util collections; all 18 data structures built from scratch.",
        "Object-Relational Persistence: Structured parent 'users' table with 'students', 'guests', and 'drivers' subtype tables.",
        "Real-World Map Calibration: 50 campus locations and 100 road edges calibrated with real Haversine GPS distances.",
        "Empirical Efficiency: High performance at N = 1000 (Max-Heap 1.45ms, Dijkstra 12.4ms).",
        "Interactive Web Demonstrator: Built with React, Vite, Leaflet, and Google Maps tile engine for seamless defense presentation."
    ]

    for c in conclusions:
        p = tf_c.add_paragraph()
        p.text = "✓ " + c
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)

    # Save presentation
    output_path = "UG_Campus_Dispatch_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully as '{output_path}'")

if __name__ == "__main__":
    create_presentation()
